import streamlit as st
from typing import List, BinaryIO
import io

# PDF parsing
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    try:
        import pdfplumber
        PDF_AVAILABLE = True
        USE_PDFPLUMBER = True
    except ImportError:
        PDF_AVAILABLE = False
        USE_PDFPLUMBER = False

# PPTX parsing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class DocumentParser:
    def __init__(self):
        """Initialize document parser with available libraries"""
        if not PDF_AVAILABLE:
            st.warning("PDF parsing not available. Install PyPDF2 or pdfplumber: pip install PyPDF2 pdfplumber")
        if not PPTX_AVAILABLE:
            st.warning("PPTX parsing not available. Install python-pptx: pip install python-pptx")
    
    def parse_document(self, uploaded_file) -> List[str]:
        """Parse uploaded document and return list of page contents"""
        
        file_extension = uploaded_file.name.lower().split('.')[-1]
        
        if file_extension == 'pdf':
            return self._parse_pdf(uploaded_file)
        elif file_extension == 'pptx':
            return self._parse_pptx(uploaded_file)
        else:
            st.error(f"Unsupported file type: {file_extension}")
            return []
    
    def _parse_pdf(self, uploaded_file) -> List[str]:
        """Parse PDF file and extract text from each page"""
        
        if not PDF_AVAILABLE:
            st.error("PDF parsing libraries not available")
            return []
        
        try:
            # Read file content
            file_content = uploaded_file.read()
            uploaded_file.seek(0)  # Reset file pointer
            
            pages = []
            
            if 'USE_PDFPLUMBER' in globals() and USE_PDFPLUMBER:
                # Use pdfplumber for better text extraction
                import pdfplumber
                
                with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        try:
                            text = page.extract_text()
                            
                            # Also try to extract tables
                            tables = page.extract_tables()
                            if tables:
                                text += "\n\nTABLES:\n"
                                for table in tables:
                                    for row in table:
                                        if row:
                                            text += " | ".join([cell or "" for cell in row]) + "\n"
                            
                            pages.append(text or f"[Page {page_num + 1} - No text extracted]")
                        except Exception as e:
                            pages.append(f"[Page {page_num + 1} - Error extracting text: {str(e)}]")
            
            else:
                # Use PyPDF2 as fallback
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                
                for page_num in range(len(pdf_reader.pages)):
                    try:
                        page = pdf_reader.pages[page_num]
                        text = page.extract_text()
                        pages.append(text or f"[Page {page_num + 1} - No text extracted]")
                    except Exception as e:
                        pages.append(f"[Page {page_num + 1} - Error extracting text: {str(e)}]")
            
            return pages
            
        except Exception as e:
            st.error(f"Error parsing PDF: {str(e)}")
            return []
    
    def _parse_pptx(self, uploaded_file) -> List[str]:
        """Parse PPTX file and extract content from each slide"""
        
        if not PPTX_AVAILABLE:
            st.error("PPTX parsing library not available")
            return []
        
        try:
            # Read file content
            file_content = uploaded_file.read()
            uploaded_file.seek(0)  # Reset file pointer
            
            prs = Presentation(io.BytesIO(file_content))
            slides = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"SLIDE {slide_num + 1}:\n"
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += f"{shape.text.strip()}\n"
                    
                    # Handle tables
                    if shape.has_table:
                        slide_text += "\nTABLE:\n"
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            slide_text += f"{row_text}\n"
                    
                    # Handle charts (basic text extraction)
                    if shape.has_chart:
                        slide_text += f"\nCHART: {shape.chart.chart_title.text_frame.text if shape.chart.chart_title else 'Untitled Chart'}\n"
                
                slides.append(slide_text)
            
            return slides
            
        except Exception as e:
            st.error(f"Error parsing PPTX: {str(e)}")
            return []
